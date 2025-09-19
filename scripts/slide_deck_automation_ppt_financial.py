#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "pandas",
#     "matplotlib",
#     "seaborn",
#     "python-pptx",
#     "openpyxl"
# ]
# ///

import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from matplotlib.dates import DateFormatter
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import warnings

warnings.filterwarnings("ignore")

input_folder = "input"
charts_folder = os.path.join(input_folder, "charts")
os.makedirs(charts_folder, exist_ok=True)

prs = Presentation()
title_slide_layout = prs.slide_layouts[5]

excel_files = [f for f in os.listdir(input_folder) if f.endswith(".xlsx")]

for file in excel_files:
    try:
        filepath = os.path.join(input_folder, file)
        df = pd.read_excel(filepath, sheet_name=0, usecols="A:Q")

        df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
        df.dropna(subset=["Date", "Stock Symbol", "Close Price", "Moving Average (50-day)",
                          "RSI (Relative Strength Index)", "MACD", "Signal Line"], inplace=True)
        df.sort_values("Date", inplace=True)

        grouped = df.groupby("Stock Symbol")

        for symbol, group in grouped:
            if group.empty:
                continue

            fig, ax1 = plt.subplots(figsize=(10, 6))
            sns.lineplot(data=group, x="Date", y="Close Price", ax=ax1, label="Close Price")
            sns.lineplot(data=group, x="Date", y="Moving Average (50-day)", ax=ax1, label="50-day MA")
            sns.lineplot(data=group, x="Date", y="RSI (Relative Strength Index)", ax=ax1, label="RSI")

            ax1.set_title(f"{symbol} - {os.path.splitext(file)[0]}", fontsize=14)
            ax1.set_xlabel("Date")
            ax1.set_ylabel("Price / RSI")
            ax1.legend(loc="upper left")
            ax1.xaxis.set_major_formatter(DateFormatter("%Y-%m-%d"))
            plt.xticks(rotation=45)

            ax2 = ax1.twinx()
            sns.lineplot(data=group, x="Date", y="MACD", ax=ax2, label="MACD", linestyle="--")
            sns.lineplot(data=group, x="Date", y="Signal Line", ax=ax2, label="Signal Line", linestyle="--")
            ax2.set_ylabel("MACD / Signal")

            chart_filename = f"{symbol}_{os.path.splitext(file)[0]}.png"
            chart_path = os.path.join(charts_folder, chart_filename)
            plt.tight_layout()
            plt.savefig(chart_path)
            plt.close()

            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = f"{os.path.splitext(file)[0]} - {symbol}"
            slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8.5))

    except Exception as e:
        print(f"Error processing {file}: {e}")

output_path = os.path.join(".", "stock_analysis.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")