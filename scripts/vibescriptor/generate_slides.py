"""
/// This script generates PowerPoint slides using OpenAI's GPT models (e.g., GPT-4) given an input topic and context.
/// Make sure to install `openai` and `python-pptx` libraries before running the script:
/// pip install openai python-pptx

import openai
from pptx import Presentation
from pptx.util import Inches, Pt

# OpenAI API setup
openai.api_key = 'your-api-key-here'  # Ensure your API key is kept secure

# Function to interact with OpenAI's GPT model
def fetch_slide_content(topic, context):
    response = openai.Completion.create(
        model="text-davinci-003",  # change to the appropriate model name
        prompt=f"Create a set of bullet points for a presentation on {topic}. Context: {context}",
        max_tokens=150,
        n=1,
        stop=None,
        temperature=0.7
    )
    return response.choices[0].text.strip()

# Function to create PowerPoint slides
def create_ppt_from_content(title, content):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(6)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Add content to the slide
    for line in content.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)

    return prs

# Main function to generate slides
def generate_presentation(topic, context):
    content = fetch_slide_content(topic, context)
    prs = create_ppt_from_content(topic, content)
    prs.save(f"{topic}.pptx")
    print(f"Presentation saved as {topic}.pptx")

# Example usage
# generate_presentation("AI in Education", "The role of Artificial Intelligence in modern education systems.")
"""
