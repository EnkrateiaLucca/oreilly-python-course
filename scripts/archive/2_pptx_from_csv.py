#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "pandas",
#     "matplotlib",
#     "python-pptx",
#     "pydantic",
#     "openai"
# ]
# ///

"""
AI-powered CSV -> PPTX report generator using:
- OpenAI Responses API with Pydantic structured outputs
- Matplotlib for charts
- python-pptx for PowerPoint generation

Usage:
    uv run ai_csv_to_ppt.py path/to/data.csv \
        --output-pptx report.pptx \
        --charts-dir charts

Environment:
    export OPENAI_API_KEY=...
"""

import argparse
import os
from dataclasses import dataclass
from typing import List, Optional, Literal

import matplotlib
matplotlib.use("Agg")  # non-interactive backend
import matplotlib.pyplot as plt
import pandas as pd
from openai import OpenAI
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches

# ---------- OpenAI / model config ----------

MODEL = "gpt-4.1-mini"  # change if you prefer another Responses-capable model


# ---------- Pydantic models for structured outputs ----------

ChartType = Literal["line", "bar", "scatter", "histogram"]


class ChartSpec(BaseModel):
    """Specification for a single chart."""
    title: str = Field(description="Title of the chart as shown in the slide.")
    description: str = Field(
        description="Short natural-language description of what this chart shows."
    )
    chart_type: ChartType = Field(
        description='One of: "line", "bar", "scatter", "histogram".'
    )
    x_column: str = Field(
        description="Column name to use on the X-axis (time, category, or index-like)."
    )
    y_columns: List[str] = Field(
        description="One or more numeric columns to plot on the Y-axis."
    )
    # Optional: simple filter to focus the data (not required)
    filter_hint: Optional[str] = Field(
        default=None,
        description=(
            "Optional human-readable filter for the subset of rows to plot, "
            "e.g. 'only year >= 2020' or 'Category == A'. This is advisory; "
            "the Python code may ignore it if it cannot be applied safely."
        ),
    )


class SlideSpec(BaseModel):
    """A single slide with bullets and optional charts."""
    title: str = Field(description="Slide title.")
    bullets: List[str] = Field(
        description="2–6 bullet points summarizing the key insight for this slide."
    )
    charts: List[ChartSpec] = Field(
        default_factory=list,
        description="0–2 charts to include on this slide.",
    )


class PresentationPlan(BaseModel):
    """Overall presentation plan."""
    title: str = Field(description="Overall presentation title.")
    subtitle: Optional[str] = Field(
        default=None, description="Optional subtitle or context."
    )
    slides: List[SlideSpec] = Field(
        description="Ordered list of slides to include."
    )


# ---------- Helper: summarize dataframe for the model ----------

def summarize_dataframe_for_llm(df: pd.DataFrame, csv_path: str, max_cols: int = 12) -> str:
    """Create a compact textual summary of the dataframe for the LLM."""
    lines: List[str] = []
    lines.append(f"Source file: {os.path.basename(csv_path)}")
    lines.append(f"Shape: {df.shape[0]} rows x {df.shape[1]} columns")
    lines.append("")

    # Basic dtype info and sample values
    cols = list(df.columns)
    if len(cols) > max_cols:
        cols = cols[:max_cols]
        lines.append(
            f"Only the first {max_cols} columns are shown below (dataset has more columns)."
        )
        lines.append("")

    for col in cols:
        series = df[col]
        dtype = str(series.dtype)
        non_null = series.dropna()
        sample_values = non_null.head(5).tolist()

        lines.append(f"Column: {col}")
        lines.append(f"  dtype: {dtype}")
        lines.append(f"  non-null count: {non_null.shape[0]}")
        if sample_values:
            lines.append("  sample values:")
            for v in sample_values:
                lines.append(f"    - {repr(v)}")
        lines.append("")

    # Basic numeric summary
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    if numeric_cols:
        lines.append("Numeric columns summary (pandas describe):")
        lines.append(df[numeric_cols].describe().to_string())
        lines.append("")

    return "\n".join(lines)


# ---------- LLM: design the presentation plan ----------

def generate_presentation_plan(client: OpenAI, df: pd.DataFrame, csv_path: str) -> PresentationPlan:
    """
    Call the OpenAI Responses API with a Pydantic structured output
    to design a presentation for the given dataframe.
    """
    dataset_summary = summarize_dataframe_for_llm(df, csv_path)

    instructions = """
You are a senior data storyteller and presentation designer.

You will receive a description of a tabular dataset (columns, dtypes, sample values, summary stats).
Your task is to design a concise, insight-focused PowerPoint-style presentation.

OUTPUT FORMAT (STRICT):
- You MUST output a JSON object that matches the provided schema:
  PresentationPlan -> slides: list[SlideSpec] -> charts: list[ChartSpec].
- Use ONLY the column names exactly as provided.
- chart_type must be one of: "line", "bar", "scatter", "histogram".
- x_column must be a single column from the dataset.
- y_columns must be a list of one or more numeric columns.

CONTENT GUIDELINES:
- Use at most 6 slides (excluding the title slide).
- Each slide:
  - 2 to 6 clear, non-redundant bullet points.
  - At most 2 charts (0, 1, or 2).
- Prefer:
  - time-like or ordered columns (e.g., dates) for x_column in line charts.
  - categorical columns for x_column in bar charts.
  - pairs of numeric columns for scatter charts.
  - single numeric column for histograms.
- Avoid:
  - using the same chart repeatedly with trivial changes.
  - referencing columns that do not exist.
- The overall title and subtitle should describe the whole dataset and its purpose (if clear).
"""

    # Structured Outputs via Responses API + Pydantic
    # See: https://platform.openai.com/docs/guides/structured-outputs
    response = client.responses.parse(
        model=MODEL,
        instructions=instructions,
        input=dataset_summary,
        temperature=0.3,
        text_format=PresentationPlan,  # Pydantic model schema
    )

    plan: PresentationPlan = response.output_parsed
    return plan


# ---------- Chart generation from ChartSpec ----------

@dataclass
class ChartContext:
    path: str
    spec: ChartSpec


def safe_apply_filter_hint(df: pd.DataFrame, filter_hint: Optional[str]) -> pd.DataFrame:
    """
    Very conservative: currently just returns the original df.

    You could extend this to parse simple hints like
    'only year >= 2020' or 'Category == A', but doing that
    robustly requires additional logic.

    For safety and generic use, we ignore filter_hint here.
    """
    _ = filter_hint  # unused for now
    return df


def create_chart_from_spec(
    df: pd.DataFrame,
    spec: ChartSpec,
    charts_dir: str,
    slide_idx: int,
    chart_idx: int,
) -> ChartContext:
    """
    Given the dataframe and a ChartSpec, generate a PNG chart
    and return its file path + context.
    """
    os.makedirs(charts_dir, exist_ok=True)

    # Apply filter (currently no-op)
    data = safe_apply_filter_hint(df, spec.filter_hint)

    # Ensure columns exist
    if spec.x_column not in data.columns:
        raise ValueError(f"x_column '{spec.x_column}' not found in dataframe.")
    for col in spec.y_columns:
        if col not in data.columns:
            raise ValueError(f"y_column '{col}' not found in dataframe.")

    x = data[spec.x_column]
    y_cols = spec.y_columns

    plt.figure(figsize=(8, 4.5))

    if spec.chart_type == "line":
        for col in y_cols:
            plt.plot(x, data[col], label=col)
    elif spec.chart_type == "bar":
        # For multi-series bar, we do grouped bars; for 1 series, simple bar
        import numpy as np

        indices = np.arange(len(x))
        width = 0.8 / max(len(y_cols), 1)
        for i, col in enumerate(y_cols):
            plt.bar(indices + i * width, data[col], width=width, label=col)
        plt.xticks(indices + (len(y_cols) - 1) * width / 2, x, rotation=45, ha="right")
    elif spec.chart_type == "scatter":
        # scatter uses first y as main; others as additional series
        for col in y_cols:
            plt.scatter(x, data[col], label=col, alpha=0.7)
    elif spec.chart_type == "histogram":
        # histogram uses first y-column only
        col = y_cols[0]
        plt.hist(data[col].dropna(), bins=30)
    else:
        raise ValueError(f"Unsupported chart_type: {spec.chart_type}")

    plt.title(spec.title)
    plt.xlabel(spec.x_column)
    if spec.chart_type != "histogram":
        plt.ylabel(", ".join(y_cols))
    else:
        plt.ylabel("Frequency")
    if len(y_cols) > 1 and spec.chart_type != "histogram":
        plt.legend()
    plt.tight_layout()

    filename = f"slide{slide_idx:02d}_chart{chart_idx:02d}.png"
    path = os.path.join(charts_dir, filename)
    plt.savefig(path)
    plt.close()

    return ChartContext(path=path, spec=spec)


# ---------- PPTX generation ----------

def build_presentation(
    plan: PresentationPlan,
    df: pd.DataFrame,
    charts_dir: str,
    csv_path: str,
    output_pptx: str,
) -> None:
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = plan.title or "AI-Generated Data Report"

    if len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_text_parts = []
        if plan.subtitle:
            subtitle_text_parts.append(plan.subtitle)
        subtitle_text_parts.append(f"Source: {os.path.basename(csv_path)}")
        subtitle_shape.text = " | ".join(subtitle_text_parts)

    # Content slides
    content_layout = prs.slide_layouts[1]  # title + content

    for i, slide_spec in enumerate(plan.slides, start=1):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_spec.title

        # Bullet text in main placeholder
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        if slide_spec.bullets:
            # First bullet as main paragraph
            p = tf.paragraphs[0]
            p.text = slide_spec.bullets[0]
            p.level = 0

            for bullet in slide_spec.bullets[1:]:
                bp = tf.add_paragraph()
                bp.text = bullet
                bp.level = 0

        # Generate and place charts
        for j, chart_spec in enumerate(slide_spec.charts[:2], start=1):  # max 2 charts visually
            try:
                ctx = create_chart_from_spec(df, chart_spec, charts_dir, i, j)
            except Exception as e:
                # If chart creation fails, add a note to the slide instead of breaking everything
                err_par = tf.add_paragraph()
                err_par.text = f"(Chart '{chart_spec.title}' could not be created: {e})"
                err_par.level = 1
                continue

            # Add the chart image to the slide
            # Simple layout: charts below bullets, from left to right
            left = Inches(0.5 + (j - 1) * 4.8)
            top = Inches(3.0)
            width = Inches(4.5)
            slide.shapes.add_picture(ctx.path, left, top, width=width)

    prs.save(output_pptx)


# ---------- Optional extra AI call: refine title/subtitle (example of multi-call usage) ----------

def refine_title_with_ai(client: OpenAI, plan: PresentationPlan, csv_path: str) -> PresentationPlan:
    """
    Second call to AI: rewrite title & subtitle to be a bit more polished.
    If it fails for any reason, return the original plan.
    """

    class TitleRefinement(BaseModel):
        title: str
        subtitle: Optional[str] = None

    try:
        instructions = """
You are a presentation title editor.
Given a draft title and subtitle plus a file name, return a polished title and subtitle.
Keep it concise and professional.
"""
        input_text = f"""
Draft title: {plan.title!r}
Draft subtitle: {plan.subtitle!r}
File name: {os.path.basename(csv_path)!r}
"""

        resp = client.responses.parse(
            model=MODEL,
            instructions=instructions,
            input=input_text,
            temperature=0.2,
            text_format=TitleRefinement,
        )
        refined = resp.output_parsed
        plan.title = refined.title
        plan.subtitle = refined.subtitle
    except Exception:
        # swallow errors; keep original plan
        pass

    return plan


# ---------- CLI / main ----------

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate an AI-designed PowerPoint report from a CSV file."
    )
    parser.add_argument("csv_path", help="Path to the input .csv file.")
    parser.add_argument(
        "--output-pptx",
        default="ai_data_report.pptx",
        help="Output PowerPoint path (default: ai_data_report.pptx).",
    )
    parser.add_argument(
        "--charts-dir",
        default="charts",
        help="Directory to store generated chart images (default: charts).",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    if not os.path.exists(args.csv_path):
        raise FileNotFoundError(f"CSV file not found: {args.csv_path}")

    print(f"Loading CSV: {args.csv_path}")
    df = pd.read_csv(args.csv_path)

    client = OpenAI()

    print("Calling OpenAI to design presentation plan (structured output)...")
    plan = generate_presentation_plan(client, df, args.csv_path)

    print("Refining title & subtitle with a second AI call...")
    plan = refine_title_with_ai(client, plan, args.csv_path)

    print("Building PowerPoint and generating charts...")
    build_presentation(
        plan=plan,
        df=df,
        charts_dir=args.charts_dir,
        csv_path=args.csv_path,
        output_pptx=args.output_pptx,
    )

    print(f"Done. Presentation saved to {args.output_pptx}")
    print(f"Charts saved under directory: {args.charts_dir}")


if __name__ == "__main__":
    main()