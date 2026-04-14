# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=0.6.23",
# ]
# ///

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Colors from brand guidelines
BLACK = RGBColor(0,0,0)
WHITE = RGBColor(255,255,255)
L1 = RGBColor(0xF9,0xF9,0xF9)
L2 = RGBColor(0xF5,0xF5,0xF5)
L3 = RGBColor(0xF0,0xF0,0xF0)
D1 = RGBColor(0x33,0x33,0x33)
D2 = RGBColor(0x55,0x55,0x55)
D3 = RGBColor(0x66,0x66,0x66)
SUCCESS_BG = RGBColor(0xE8,0xF5,0xE8)
WARN_BG = RGBColor(0xFF,0xEE,0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 default wide
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer(slide):
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12.33)
    height = Inches(0.3)
    tx = slide.shapes.add_textbox(left, top, width, height)
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "© 2025 Automata Learning Lab | Professional educational materials crafted with care"
    run.font.name = "Helvetica"
    run.font.size = Pt(8)
    run.font.color.rgb = D3
    p.alignment = PP_ALIGN.CENTER

def add_title(slide, text, subtitle=None):
    # big uppercase title
    left = Inches(0.7); top = Inches(0.6); width = Inches(11.9); height = Inches(1.2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Helvetica"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Helvetica"
        p2.font.size = Pt(18)
        p2.font.color.rgb = D2
        p2.space_before = Pt(10)

def add_bullets(slide, bullets, left=0.9, top=1.8, width=11.6, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Helvetica"
        p.font.size = Pt(font_size)
        p.font.color.rgb = D1
        p.space_after = Pt(6)

def add_card(slide, title, body_lines, left, top, width, height, bg=L1, border=BLACK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Helvetica"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.space_after = Pt(4)
    for line in body_lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.level = 0
        bp.font.name = "Helvetica"
        bp.font.size = Pt(13)
        bp.font.color.rgb = D1

def section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title(slide, title)
    add_footer(slide)
    return slide

# --- Slides content ---

# 1. Title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Context Engineering 2.0", "The Context of Context Engineering (SII-GAIR, 2025)")
# author line
box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1))
tf = box.text_frame
p=tf.paragraphs[0]
p.text="Key ideas & design considerations for modern AI agents"
p.font.name="Helvetica"; p.font.size=Pt(20); p.font.color.rgb=D2
add_footer(slide)

# 2. Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Agenda")
add_bullets(slide, [
    "Why context engineering matters now",
    "Entropy reduction view + formal definition",
    "Historical evolution: eras 1.0 → 4.0",
    "Design considerations: collect, manage, use",
    "Applications, challenges, future directions"
])
add_footer(slide)

# 3. Motivation / problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Motivation")
add_bullets(slide, [
    "LLMs are sensitive to what’s inside the context window.",
    "Long‑horizon reasoning needs more than a bigger window.",
    "We need systematic ways to align machine behavior with human intent."
])
add_footer(slide)

# 4. What is context engineering?
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "What is Context Engineering?")
add_bullets(slide, [
    "Practice of designing, organizing, and managing context so machines act on human intentions.",
    "Includes prompting, RAG, tool use, memory, multimodal inputs, and context sharing.",
])
# callout
add_card(slide, "Core idea = Entropy Reduction", [
    "Humans fill gaps implicitly; machines can’t.",
    "So we compress high‑entropy real‑world intent",
    "into low‑entropy representations a model can use."
], left=0.9, top=3.3, width=11.5, height=2.2, bg=L2)
add_footer(slide)

# 5. Historical perspective
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "A Longer History")
add_bullets(slide, [
    "Context engineering predates LLMs by ~20 years.",
    "Roots in ubiquitous computing, context‑aware systems, HCI.",
    "Each leap in machine intelligence changes how we build interfaces."
])
add_footer(slide)

# 6. Four-stage trajectory
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Four Eras of Context Engineering")
add_card(slide, "1.0 Primitive Computation", [
    "1990s–2020",
    "Structured, low‑entropy inputs",
    "Humans translate intent for machines"
], left=0.7, top=1.9, width=3.1, height=2.2, bg=L1)
add_card(slide, "2.0 Agent‑Centric Intelligence", [
    "2020–present",
    "Natural language + ambiguity tolerance",
    "Prompting, RAG, tools, memory"
], left=3.95, top=1.9, width=3.3, height=2.2, bg=L1)
add_card(slide, "3.0 Human‑Level (future)", [
    "Human‑like sensing & understanding",
    "Seamless collaboration",
    "Less explicit context mgmt"
], left=7.45, top=1.9, width=2.9, height=2.2, bg=L1)
add_card(slide, "4.0 Superhuman (speculative)", [
    "Models construct context proactively",
    "Reveal hidden needs",
    "‘God’s‑eye view’ of intent"
], left=10.55, top=1.9, width=2.1, height=2.2, bg=L1)
add_bullets(slide, ["Trend: increasing intelligence → decreasing human effort"], top=4.7, font_size=18)
add_footer(slide)

# 7. Design considerations overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Design Considerations (Era 2.0)")
add_bullets(slide, [
    "Three dimensions form the pipeline:",
    "1) Context Collection & Storage",
    "2) Context Management",
    "3) Context Usage"
], top=1.9)
add_footer(slide)

# 8. Collection principles
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Context Collection & Storage")
add_bullets(slide, [
    "Multimodal sources: text, images, audio, sensors, logs, APIs.",
    "Distributed storage: on‑device, edge, cloud, layered by latency.",
])
add_card(slide, "Principle: Minimal Sufficiency", [
    "Collect/store only what's needed for the task.",
    "Value is sufficiency, not volume."
], left=0.9, top=3.2, width=5.8, height=2.2, bg=SUCCESS_BG)
add_card(slide, "Principle: Semantic Continuity", [
    "Preserve continuity of meaning, not just raw data.",
    "Maintain coherence across time and devices."
], left=6.6, top=3.2, width=5.6, height=2.2, bg=L2)
add_footer(slide)

# 9. Context management
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Context Management")
add_bullets(slide, [
    "Process raw context: clean, normalize, align modalities.",
    "Organize memory in layers: short‑term → long‑term.",
    "Isolate context to avoid contamination across tasks.",
    "Abstract/compress using summaries or embeddings (‘self‑baking’)."
])
add_footer(slide)

# 10. Context usage
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Context Usage")
add_bullets(slide, [
    "Intra‑system sharing: agents pass context via prompts, schemas, or shared memory.",
    "Cross‑system sharing: adapters or common representations between tools/agents.",
    "Context selection: pick relevant pieces for understanding and action.",
    "Proactive inference: mine preferences, infer hidden goals, offer help.",
    "Lifelong preservation: update personal context without drift."
])
add_footer(slide)

# 11. Application: Gemini CLI
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Applications: Coding CLIs")
add_bullets(slide, [
    "Example: Gemini CLI uses GEMINI.md files as project context.",
    "Static context loads at startup; dynamic context accumulates in dialog.",
    "History is periodically summarized into compact ‘reasoning state’."
])
add_footer(slide)

# 12. Application: Deep Research
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Applications: Deep Research Agents")
add_bullets(slide, [
    "Open‑ended, long‑horizon search + reasoning.",
    "Need periodic compression to stay within window.",
    "Context snapshots preserve evidence and guide next searches."
])
add_footer(slide)

# 13. Application: BCI
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Applications: Brain‑Computer Interfaces")
add_bullets(slide, [
    "BCIs can collect internal user state directly (attention, emotion, cognitive load).",
    "Richer + more convenient context collection beyond language.",
    "Still early: noisy signals, coarse interpretation."
])
add_footer(slide)

# 14. Challenges
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Challenges & Future Directions")
add_bullets(slide, [
    "Collection is still inefficient; users can’t always articulate intent.",
    "Storage bottlenecks at lifelong scale.",
    "Processing degradation in long contexts (latency + reasoning quality).",
    "System instability and semantic drift over time.",
    "Need better architectures + adaptive selection + evaluation frameworks."
])
add_footer(slide)

# 15. Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Conclusion")
add_bullets(slide, [
    "Context engineering is a long‑evolving discipline.",
    "Core job: reduce entropy between human intent and machine action.",
    "Era 2.0 requires systematic pipelines for collection, management, and usage.",
    "As models improve, human effort shifts from explicit context handling to oversight."
])
add_footer(slide)

# 16. Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Q&A")
add_footer(slide)

out_path = "./context-engineering-2.0-presentation.pptx"
prs.save(out_path)
out_path