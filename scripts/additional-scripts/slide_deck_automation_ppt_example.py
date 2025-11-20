# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-pptx",
#   "matplotlib",
#   "pillow",
# ]
# ///
"""
Generates a 5-slide PPTX on how LLMs augment productivity.
Slides:
1) Title
2) Chart (LLM usage/adoption) with sources
3) Diagram: How LLMs improve productivity (shapes + arrows)
4) Creative: “Before vs After” workflow visual + callout
5) Conclusion
Data sources (fetched & verified by you earlier; cited on slides):
- Stack Overflow Dev Survey 2024 (AI tool usage 62% currently use / 76% use-or-plan)
  https://survey.stackoverflow.co/2024/ai
- McKinsey 'State of AI' 2024 (65% regularly using genAI)
  https://www.mckinsey.com/capabilities/quantumblack/our-insights/the-state-of-ai-2024
- Pew Research (Jun 25, 2025): 34% of U.S. adults have used ChatGPT
  https://www.pewresearch.org/short-reads/2025/06/25/34-of-us-adults-have-used-chatgpt-about-double-the-share-in-2023/
- VS Magazine summary of a Copilot study: ~55% faster on a coding task
  https://visualstudiomagazine.com/articles/2024/09/17/another-report-weighs-in-on-github-copilot-dev-productivity.aspx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement

import io
import matplotlib.pyplot as plt

def add_metadata(slide, key, value):
    # tiny hidden textbox for metadata
    tx = slide.shapes.add_textbox(Inches(0.1), Inches(6.9), Inches(2.5), Inches(0.25))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{key}: {value}"
    run.font.size = Pt(1)
    run.font.color.rgb = RGBColor(255, 255, 255)

def add_footnote(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(100, 100, 100)

def add_title(prs, title_text, subtitle_text=None, layout_idx=0):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    if subtitle_text is not None and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide

def add_bullets(slide, title, bullets):
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0

def add_matplotlib_chart_image(slide, fig, left_in=0.7, top_in=1.8, width_in=8.0):
    bio = io.BytesIO()
    fig.savefig(bio, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    bio.seek(0)
    slide.shapes.add_picture(bio, Inches(left_in), Inches(top_in), width=Inches(width_in))

def make_usage_chart():
    """
    Bars comparing different, reputable metrics (different populations):
    - Developers using AI tools (StackOverflow 2024): 62%
    - Organizations regularly using GenAI (McKinsey 2024): 65%
    - U.S. adults who have ever used ChatGPT (Pew 2025): 34%
    """
    labels = [
        "Devs using AI tools (2024)",
        "Orgs using GenAI (2024)",
        "U.S. adults used ChatGPT (2025)",
    ]
    values = [62, 65, 34]

    fig = plt.figure(figsize=(8, 3.8))
    ax = fig.add_subplot(111)
    bars = ax.bar(labels, values)
    ax.set_ylim(0, 100)
    ax.set_ylabel("Percent")
    ax.set_title("LLM/GenAI Adoption Snapshot (Different Populations)")
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width()/2, v + 2, f"{v}%", ha="center", va="bottom", fontsize=10)
    fig.tight_layout()
    return fig

def add_rectangle(slide, x, y, w, h, text, font_size=14, fill=(240,240,240)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fillfmt = shape.fill
    fillfmt.solid()
    fillfmt.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(font_size)
    return shape

def add_arrow(slide, x1, y1, x2, y2):
    # simple straight arrow: draw a line; pptx has connectors but this is robust
    line = slide.shapes.add_connector(
        MSO_LINE.SIMPLE, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.end_arrowhead = True
    line.line.width = Pt(2.25)
    line.line.color.rgb = RGBColor(120, 120, 120)
    return line

def create_pptx(path="llm_productivity.pptx"):
    prs = Presentation()

    # --- Slide 1: Title ---
    slide1 = add_title(
        prs,
        "How Large Language Models (LLMs) Augment Productivity",
        "Usage • Mechanisms • Impact"
    )
    add_metadata(slide1, "slide_number", "1")
    add_metadata(slide1, "type", "title")

    # --- Slide 2: Chart (verified sources) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title + Content
    slide2.shapes.title.text = "LLM/GenAI Adoption (Selected Metrics)"
    # chart
    fig = make_usage_chart()
    add_matplotlib_chart_image(slide2, fig)
    # footnotes / sources
    add_footnote(
        slide2,
        "Sources: Stack Overflow Dev Survey 2024 (AI tools use 62%), "
        "McKinsey State of AI 2024 (65% genAI use), Pew Research 2025 (34% U.S. adults used ChatGPT)."
    )
    add_metadata(slide2, "slide_number", "2")
    add_metadata(slide2, "type", "chart")

    # --- Slide 3: Diagram — How LLMs augment productivity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Mechanisms: How LLMs Boost Productivity"

    # boxes
    b1 = add_rectangle(slide3, 0.6, 1.6, 2.6, 1.0, "Inputs\n(Data • Prompts • Context)", 14)
    b2 = add_rectangle(slide3, 3.6, 1.6, 3.0, 1.0, "LLM\n(Reason • Generate • Retrieve)", 14, fill=(225,240,255))
    b3 = add_rectangle(slide3, 7.0, 1.6, 2.8, 1.0, "Outputs\n(Drafts • Code • Summaries)", 14)
    add_arrow(slide3, 3.2, 2.1, 3.55, 2.1)
    add_arrow(slide3, 6.6, 2.1, 6.95, 2.1)

    # bottom row: feedback loop
    b4 = add_rectangle(slide3, 1.4, 3.2, 2.4, 0.9, "Human Review\n& Edit", 12, fill=(245,235,255))
    b5 = add_rectangle(slide3, 4.2, 3.2, 2.4, 0.9, "Tooling\n(Checks • Tests • Search)", 12, fill=(245,255,235))
    b6 = add_rectangle(slide3, 6.9, 3.2, 2.4, 0.9, "Refine\n(Iterate • Re-prompt)", 12, fill=(235,245,255))
    add_arrow(slide3, 2.6, 2.6, 2.6, 3.15)
    add_arrow(slide3, 5.4, 2.6, 5.4, 3.15)
    add_arrow(slide3, 8.1, 2.6, 8.1, 3.15)
    add_arrow(slide3, 3.8, 3.65, 4.15, 3.65)
    add_arrow(slide3, 6.6, 3.65, 6.95, 3.65)
    add_arrow(slide3, 8.1, 3.65, 8.1, 2.6)  # back up to outputs (loop)

    add_footnote(slide3, "Flow: Inputs → LLM → Outputs, with human & tooling feedback loops to improve quality and speed.")
    add_metadata(slide3, "slide_number", "3")
    add_metadata(slide3, "type", "diagram")

    # --- Slide 4: Creative slide (visual 'Before vs After' + callout) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Creative View: From Tasks to Leverage"

    # Before / After panels
    before = add_rectangle(slide4, 0.6, 1.6, 4.2, 3.4,
                           "Before LLMs\n\n• Manual drafting & research\n• Repetitive formatting\n• Context switching\n• Slow iteration",
                           14, fill=(255,235,235))
    after = add_rectangle(slide4, 5.1, 1.6, 4.7, 3.4,
                          "With LLMs as Copilots\n\n• Drafts & summaries in minutes\n• Assisted research & retrieval\n• Inline code/help\n• Faster iterations",
                          14, fill=(235,255,240))
    add_arrow(slide4, 4.9, 3.3, 5.05, 3.3)

    # Big callout with one quantified stat
    callout = slide4.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(2.0), Inches(5.3), Inches(6.8), Inches(1.2))
    tf = callout.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Reported speed-ups: ~55% faster on a coding task (Copilot study summary)"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(16)
    callout.line.color.rgb = RGBColor(200, 200, 200)
    add_footnote(slide4, "Source: Visual Studio Magazine (Sep 17, 2024) summarizing a GitHub Copilot study.")
    add_metadata(slide4, "slide_number", "4")
    add_metadata(slide4, "type", "creative")

    # --- Slide 5: Conclusion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Conclusion"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    for line in [
        "Adoption is widespread across developers, organizations, and consumers.",
        "LLMs shift time from grunt work to higher-value thinking and iteration.",
        "Best results: combine LLMs with verification, tooling, and human judgment.",
        "Design for feedback loops, reproducibility, and responsible use."
    ]:
        p = tf5.add_paragraph() if len(tf5.paragraphs[0].text) else tf5.paragraphs[0]
        p.text = "• " + line
        p.level = 0

    add_metadata(slide5, "slide_number", "5")
    add_metadata(slide5, "type", "conclusion")

    prs.save(path)
    print(f"Saved presentation to {path}")

if __name__ == "__main__":
    create_pptx()