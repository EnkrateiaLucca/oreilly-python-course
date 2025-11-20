#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx", "anthropic", "python-dotenv", "pandas", "matplotlib", "requests", "pillow"]
# ///

"""
MM_powerpoint_generator.py
Student: MM - "powerpoint generation"

This script demonstrates automated PowerPoint presentation generation from data,
text, and various content sources using AI for content creation and formatting.

Key Learning Objectives:
- Automated presentation creation
- Data visualization integration
- AI-powered content generation
- Template-based design automation
- Multi-source content integration
"""

import os
import json
import pandas as pd
import matplotlib.pyplot as plt
from datetime import datetime
from typing import List, Dict, Any, Optional
from io import BytesIO
from dotenv import load_dotenv
import anthropic

# PowerPoint manipulation library
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Load environment variables
load_dotenv()

class PowerPointGenerator:
    """
    A class to automatically generate PowerPoint presentations from various data sources.
    This demonstrates automated slide creation, data visualization, and AI-powered content.
    """

    def __init__(self):
        """Initialize the PowerPoint generator with AI client."""
        # Initialize Anthropic client for AI-powered content generation
        self.ai_client = anthropic.Anthropic(
            api_key=os.getenv("ANTHROPIC_API_KEY")
        )

        # Define presentation themes and color schemes
        self.themes = {
            "corporate": {
                "primary_color": RGBColor(0, 51, 102),      # Dark blue
                "secondary_color": RGBColor(255, 255, 255),  # White
                "accent_color": RGBColor(0, 153, 204),       # Light blue
                "text_color": RGBColor(51, 51, 51)           # Dark gray
            },
            "modern": {
                "primary_color": RGBColor(45, 45, 45),       # Charcoal
                "secondary_color": RGBColor(255, 255, 255),  # White
                "accent_color": RGBColor(255, 87, 51),       # Orange
                "text_color": RGBColor(85, 85, 85)           # Gray
            },
            "minimal": {
                "primary_color": RGBColor(34, 34, 34),       # Dark gray
                "secondary_color": RGBColor(248, 248, 248),  # Light gray
                "accent_color": RGBColor(102, 204, 0),       # Green
                "text_color": RGBColor(68, 68, 68)           # Medium gray
            }
        }

        # Slide layout types
        self.layout_types = {
            "title": 0,
            "title_content": 1,
            "section_header": 2,
            "two_content": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6,
            "content_caption": 7,
            "picture_caption": 8
        }

    def get_sample_presentation_data(self) -> Dict[str, Any]:
        """
        Generate sample data for presentation creation.
        In a real scenario, this would come from databases, APIs, or user input.
        """
        return {
            "title": "Q4 2024 Business Performance Review",
            "subtitle": "Strategic Insights and Future Outlook",
            "author": "Data Analytics Team",
            "company": "TechCorp Solutions",
            "date": datetime.now().strftime("%B %d, %Y"),
            "theme": "corporate",
            "sections": [
                {
                    "title": "Executive Summary",
                    "type": "overview",
                    "content": {
                        "key_points": [
                            "Revenue increased 25% year-over-year",
                            "Customer satisfaction reached 92%",
                            "Expanded into 3 new international markets",
                            "Launched 2 major product features"
                        ],
                        "metrics": {
                            "revenue": "$2.4M",
                            "customers": "15,000+",
                            "growth_rate": "25%",
                            "market_share": "18%"
                        }
                    }
                },
                {
                    "title": "Financial Performance",
                    "type": "data_analysis",
                    "content": {
                        "revenue_data": {
                            "Q1": 1800000,
                            "Q2": 2100000,
                            "Q3": 2200000,
                            "Q4": 2400000
                        },
                        "expense_data": {
                            "Q1": 1200000,
                            "Q2": 1300000,
                            "Q3": 1350000,
                            "Q4": 1400000
                        },
                        "profit_margin": {
                            "Q1": 33.3,
                            "Q2": 38.1,
                            "Q3": 38.6,
                            "Q4": 41.7
                        }
                    }
                },
                {
                    "title": "Customer Analytics",
                    "type": "insights",
                    "content": {
                        "satisfaction_scores": {
                            "Product Quality": 4.5,
                            "Customer Service": 4.3,
                            "Pricing": 4.1,
                            "User Experience": 4.6,
                            "Support Response": 4.2
                        },
                        "demographics": {
                            "Enterprise": 45,
                            "Mid-Market": 35,
                            "Small Business": 20
                        }
                    }
                },
                {
                    "title": "Market Opportunities",
                    "type": "strategic",
                    "content": {
                        "opportunities": [
                            "AI integration features showing 40% user interest",
                            "Mobile app expansion with 60% customer demand",
                            "European market entry projecting $1M additional revenue",
                            "Partnership opportunities with 3 major platforms"
                        ],
                        "challenges": [
                            "Increased competition in core market",
                            "Rising development costs",
                            "Talent acquisition in AI/ML space",
                            "Regulatory compliance in new markets"
                        ]
                    }
                },
                {
                    "title": "2025 Strategic Roadmap",
                    "type": "roadmap",
                    "content": {
                        "q1_goals": [
                            "Launch AI-powered analytics dashboard",
                            "Complete European market entry",
                            "Hire 5 additional engineers"
                        ],
                        "q2_goals": [
                            "Mobile app beta release",
                            "Partnership integration completion",
                            "Customer success program expansion"
                        ],
                        "annual_targets": {
                            "revenue_goal": "$3.5M",
                            "customer_goal": "25,000",
                            "market_expansion": "2 new regions"
                        }
                    }
                }
            ]
        }

    def generate_content_with_ai(self, section_data: Dict[str, Any], context: str) -> Dict[str, Any]:
        """
        Use AI to generate and enhance presentation content.

        Args:
            section_data: Raw section data
            context: Additional context about the presentation

        Returns:
            Enhanced content with AI-generated text and insights
        """

        section_title = section_data.get("title", "Untitled Section")
        section_content = section_data.get("content", {})

        # Create AI prompt for content enhancement
        prompt = f"""
        Create engaging presentation content for a business slide titled "{section_title}".

        Context: {context}

        Section Data: {json.dumps(section_content, indent=2)}

        Generate:
        1. A compelling slide title (if different from current)
        2. 3-5 bullet points highlighting key insights
        3. A brief narrative summary (2-3 sentences)
        4. Suggested talking points for the presenter
        5. Key takeaways for the audience

        Focus on:
        - Clear, actionable insights
        - Business impact and implications
        - Data-driven conclusions
        - Professional tone suitable for executives

        Respond in JSON format:
        {{
            "enhanced_title": "slide title",
            "key_insights": ["insight1", "insight2", "insight3"],
            "narrative": "summary paragraph",
            "talking_points": ["point1", "point2", "point3"],
            "key_takeaways": ["takeaway1", "takeaway2"],
            "recommended_visuals": ["chart_type1", "chart_type2"]
        }}
        """

        try:
            # Make API call to Anthropic
            response = self.ai_client.messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=1000,
                messages=[{"role": "user", "content": prompt}]
            )

            # Parse the AI response
            ai_content = json.loads(response.content[0].text)
            return ai_content

        except Exception as e:
            print(f"AI content generation failed for {section_title}: {str(e)}")
            # Fallback to basic content
            return {
                "enhanced_title": section_title,
                "key_insights": ["Data analysis completed", "Key metrics identified", "Insights extracted"],
                "narrative": f"This section covers {section_title.lower()} with relevant data and analysis.",
                "talking_points": ["Review the data", "Discuss implications", "Plan next steps"],
                "key_takeaways": ["Important metrics tracked", "Progress measured"],
                "recommended_visuals": ["bar_chart", "line_chart"]
            }

    def create_chart_from_data(self, chart_data: Dict[str, Any], chart_type: str = "bar") -> str:
        """
        Create a chart from data and save as image file.

        Args:
            chart_data: Data for chart creation
            chart_type: Type of chart (bar, line, pie)

        Returns:
            Filename of the created chart image
        """

        plt.figure(figsize=(10, 6))
        plt.style.use('seaborn-v0_8')  # Modern chart styling

        filename = f"chart_{datetime.now().strftime('%H%M%S')}.png"

        if chart_type == "bar":
            # Create bar chart
            keys = list(chart_data.keys())
            values = list(chart_data.values())
            bars = plt.bar(keys, values, color=['#003366', '#0099CC', '#66CCFF', '#99DDFF'])

            # Customize appearance
            plt.title("Performance Metrics", fontsize=16, fontweight='bold', pad=20)
            plt.ylabel("Values", fontsize=12)
            plt.xticks(rotation=45, ha='right')

            # Add value labels on bars
            for bar, value in zip(bars, values):
                plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.01,
                        f'{value:,.0f}' if isinstance(value, (int, float)) else str(value),
                        ha='center', va='bottom', fontweight='bold')

        elif chart_type == "line":
            # Create line chart
            keys = list(chart_data.keys())
            values = list(chart_data.values())
            plt.plot(keys, values, marker='o', linewidth=3, markersize=8, color='#003366')

            plt.title("Trend Analysis", fontsize=16, fontweight='bold', pad=20)
            plt.ylabel("Values", fontsize=12)
            plt.xticks(rotation=45, ha='right')
            plt.grid(True, alpha=0.3)

        elif chart_type == "pie":
            # Create pie chart
            keys = list(chart_data.keys())
            values = list(chart_data.values())
            colors = ['#003366', '#0099CC', '#66CCFF', '#99DDFF', '#CCE6FF']

            plt.pie(values, labels=keys, autopct='%1.1f%%', startangle=90,
                   colors=colors[:len(keys)])
            plt.title("Distribution Analysis", fontsize=16, fontweight='bold', pad=20)

        plt.tight_layout()
        plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()

        return filename

    def create_presentation(self, presentation_data: Dict[str, Any]) -> str:
        """
        Create a complete PowerPoint presentation from structured data.

        Args:
            presentation_data: Complete presentation data structure

        Returns:
            Filename of the created presentation
        """

        print("🎨 Creating PowerPoint presentation...")

        # Create new presentation
        prs = Presentation()

        # Get theme colors
        theme = self.themes.get(presentation_data.get("theme", "corporate"))

        # Set presentation properties
        prs.core_properties.title = presentation_data["title"]
        prs.core_properties.author = presentation_data["author"]
        prs.core_properties.created = datetime.now()

        # Create title slide
        self._create_title_slide(prs, presentation_data, theme)

        # Create agenda slide
        self._create_agenda_slide(prs, presentation_data, theme)

        # Process each section
        for i, section in enumerate(presentation_data["sections"]):
            print(f"📊 Processing section: {section['title']}")

            # Generate AI-enhanced content
            context = f"Business presentation for {presentation_data['company']} on {presentation_data['title']}"
            enhanced_content = self.generate_content_with_ai(section, context)

            # Create section slides based on type
            if section["type"] == "overview":
                self._create_overview_slide(prs, section, enhanced_content, theme)

            elif section["type"] == "data_analysis":
                self._create_data_slide(prs, section, enhanced_content, theme)

            elif section["type"] == "insights":
                self._create_insights_slide(prs, section, enhanced_content, theme)

            elif section["type"] == "strategic":
                self._create_strategic_slide(prs, section, enhanced_content, theme)

            elif section["type"] == "roadmap":
                self._create_roadmap_slide(prs, section, enhanced_content, theme)

        # Create conclusion slide
        self._create_conclusion_slide(prs, presentation_data, theme)

        # Save presentation
        filename = f"generated_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(filename)

        print(f"✅ Presentation saved as: {filename}")
        return filename

    def _create_title_slide(self, prs: Presentation, data: Dict[str, Any], theme: Dict[str, Any]):
        """Create the title slide."""
        slide_layout = prs.slide_layouts[self.layout_types["title"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = data["title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True

        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"{data['subtitle']}\n\n{data['company']}\n{data['date']}"
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.color.rgb = theme["text_color"]
            paragraph.font.size = Pt(20)

    def _create_agenda_slide(self, prs: Presentation, data: Dict[str, Any], theme: Dict[str, Any]):
        """Create an agenda slide."""
        slide_layout = prs.slide_layouts[self.layout_types["title_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = "Agenda"
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True

        # Add agenda items
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        for i, section in enumerate(data["sections"], 1):
            p = text_frame.add_paragraph()
            p.text = f"{i}. {section['title']}"
            p.font.size = Pt(24)
            p.font.color.rgb = theme["text_color"]
            p.space_after = Pt(12)

    def _create_overview_slide(self, prs: Presentation, section: Dict[str, Any],
                             enhanced_content: Dict[str, Any], theme: Dict[str, Any]):
        """Create an overview slide with key metrics."""
        slide_layout = prs.slide_layouts[self.layout_types["two_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = enhanced_content["enhanced_title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Left content: Key insights
        left_content = slide.placeholders[1]
        text_frame = left_content.text_frame
        text_frame.clear()

        # Add heading
        p = text_frame.add_paragraph()
        p.text = "Key Insights"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = theme["accent_color"]

        # Add insights
        for insight in enhanced_content["key_insights"]:
            p = text_frame.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(16)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

        # Right content: Metrics
        right_content = slide.placeholders[2]
        metrics = section["content"].get("metrics", {})

        if metrics:
            text_frame = right_content.text_frame
            text_frame.clear()

            # Add heading
            p = text_frame.add_paragraph()
            p.text = "Key Metrics"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = theme["accent_color"]

            # Add metrics
            for metric, value in metrics.items():
                p = text_frame.add_paragraph()
                p.text = f"{metric.replace('_', ' ').title()}: {value}"
                p.font.size = Pt(18)
                p.font.color.rgb = theme["text_color"]
                p.font.bold = True
                p.level = 1

    def _create_data_slide(self, prs: Presentation, section: Dict[str, Any],
                          enhanced_content: Dict[str, Any], theme: Dict[str, Any]):
        """Create a data analysis slide with charts."""
        slide_layout = prs.slide_layouts[self.layout_types["title_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = enhanced_content["enhanced_title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Create and insert chart
        revenue_data = section["content"].get("revenue_data", {})
        if revenue_data:
            chart_filename = self.create_chart_from_data(revenue_data, "line")

            # Add chart to slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4.5)

            try:
                slide.shapes.add_picture(chart_filename, left, top, width, height)
                # Clean up chart file
                os.remove(chart_filename)
            except Exception as e:
                print(f"Could not add chart to slide: {e}")

        # Add narrative text
        left = Inches(1)
        top = Inches(6.8)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = enhanced_content["narrative"]
        p.font.size = Pt(14)
        p.font.color.rgb = theme["text_color"]

    def _create_insights_slide(self, prs: Presentation, section: Dict[str, Any],
                             enhanced_content: Dict[str, Any], theme: Dict[str, Any]):
        """Create an insights slide with customer data visualization."""
        slide_layout = prs.slide_layouts[self.layout_types["two_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = enhanced_content["enhanced_title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Left content: Key insights
        left_content = slide.placeholders[1]
        text_frame = left_content.text_frame
        text_frame.clear()

        for insight in enhanced_content["key_insights"]:
            p = text_frame.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(16)
            p.font.color.rgb = theme["text_color"]

        # Right content: Customer satisfaction scores
        satisfaction_scores = section["content"].get("satisfaction_scores", {})
        if satisfaction_scores:
            right_content = slide.placeholders[2]
            text_frame = right_content.text_frame
            text_frame.clear()

            # Add heading
            p = text_frame.add_paragraph()
            p.text = "Satisfaction Scores"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = theme["accent_color"]

            # Add scores
            for category, score in satisfaction_scores.items():
                p = text_frame.add_paragraph()
                stars = "★" * int(score) + "☆" * (5 - int(score))
                p.text = f"{category}: {score}/5 {stars}"
                p.font.size = Pt(14)
                p.font.color.rgb = theme["text_color"]
                p.level = 1

    def _create_strategic_slide(self, prs: Presentation, section: Dict[str, Any],
                              enhanced_content: Dict[str, Any], theme: Dict[str, Any]):
        """Create a strategic opportunities slide."""
        slide_layout = prs.slide_layouts[self.layout_types["two_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = enhanced_content["enhanced_title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Left content: Opportunities
        left_content = slide.placeholders[1]
        text_frame = left_content.text_frame
        text_frame.clear()

        # Add heading
        p = text_frame.add_paragraph()
        p.text = "Opportunities"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 153, 0)  # Green

        opportunities = section["content"].get("opportunities", [])
        for opportunity in opportunities:
            p = text_frame.add_paragraph()
            p.text = f"+ {opportunity}"
            p.font.size = Pt(14)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

        # Right content: Challenges
        right_content = slide.placeholders[2]
        text_frame = right_content.text_frame
        text_frame.clear()

        # Add heading
        p = text_frame.add_paragraph()
        p.text = "Challenges"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(204, 51, 0)  # Red

        challenges = section["content"].get("challenges", [])
        for challenge in challenges:
            p = text_frame.add_paragraph()
            p.text = f"- {challenge}"
            p.font.size = Pt(14)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

    def _create_roadmap_slide(self, prs: Presentation, section: Dict[str, Any],
                            enhanced_content: Dict[str, Any], theme: Dict[str, Any]):
        """Create a strategic roadmap slide."""
        slide_layout = prs.slide_layouts[self.layout_types["title_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = enhanced_content["enhanced_title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Content area
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        # Q1 Goals
        p = text_frame.add_paragraph()
        p.text = "Q1 2025 Goals"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = theme["accent_color"]

        q1_goals = section["content"].get("q1_goals", [])
        for goal in q1_goals:
            p = text_frame.add_paragraph()
            p.text = f"• {goal}"
            p.font.size = Pt(14)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

        # Q2 Goals
        p = text_frame.add_paragraph()
        p.text = "Q2 2025 Goals"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = theme["accent_color"]

        q2_goals = section["content"].get("q2_goals", [])
        for goal in q2_goals:
            p = text_frame.add_paragraph()
            p.text = f"• {goal}"
            p.font.size = Pt(14)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

        # Annual Targets
        p = text_frame.add_paragraph()
        p.text = "2025 Annual Targets"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = theme["accent_color"]

        annual_targets = section["content"].get("annual_targets", {})
        for target, value in annual_targets.items():
            p = text_frame.add_paragraph()
            p.text = f"• {target.replace('_', ' ').title()}: {value}"
            p.font.size = Pt(14)
            p.font.color.rgb = theme["text_color"]
            p.level = 1

    def _create_conclusion_slide(self, prs: Presentation, data: Dict[str, Any], theme: Dict[str, Any]):
        """Create a conclusion slide."""
        slide_layout = prs.slide_layouts[self.layout_types["title_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = "Thank You"
        title.text_frame.paragraphs[0].font.color.rgb = theme["primary_color"]
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True

        # Add contact information
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        p = text_frame.add_paragraph()
        p.text = "Questions & Discussion"
        p.font.size = Pt(24)
        p.font.color.rgb = theme["accent_color"]
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = f"\n{data['company']}\n{data['author']}\n{data['date']}"
        p.font.size = Pt(18)
        p.font.color.rgb = theme["text_color"]
        p.alignment = PP_ALIGN.CENTER

    def generate_presentation_from_template(self, template_type: str, data_source: str) -> str:
        """
        Generate a presentation from predefined templates and data sources.

        Args:
            template_type: Type of presentation template
            data_source: Source of data (json file, csv, api, etc.)

        Returns:
            Filename of the generated presentation
        """

        print(f"🎯 Generating {template_type} presentation from {data_source}")

        # For demo purposes, use sample data
        # In real implementation, this would read from actual data sources
        presentation_data = self.get_sample_presentation_data()

        # Customize based on template type
        if template_type == "sales_report":
            presentation_data["title"] = "Sales Performance Report"
            presentation_data["theme"] = "corporate"

        elif template_type == "project_update":
            presentation_data["title"] = "Project Status Update"
            presentation_data["theme"] = "modern"

        elif template_type == "quarterly_review":
            presentation_data["title"] = "Quarterly Business Review"
            presentation_data["theme"] = "minimal"

        # Generate the presentation
        return self.create_presentation(presentation_data)

def main():
    """
    Main function to demonstrate PowerPoint generation capabilities.
    """

    print("🚀 PowerPoint Generator Demo")
    print("=" * 50)

    # Initialize the generator
    generator = PowerPointGenerator()

    print("📋 Available presentation templates:")
    print("1. Sales Report")
    print("2. Project Update")
    print("3. Quarterly Review")
    print("4. Custom Data Presentation")
    print()

    # Generate sample presentation
    print("🎨 Creating sample business presentation...")
    presentation_data = generator.get_sample_presentation_data()
    filename = generator.create_presentation(presentation_data)

    print(f"✅ Presentation created successfully: {filename}")

    # Demonstrate template-based generation
    print("\n📊 Generating additional template examples...")

    # Sales report
    sales_filename = generator.generate_presentation_from_template("sales_report", "sample_data")
    print(f"📈 Sales report generated: {sales_filename}")

    # Project update
    project_filename = generator.generate_presentation_from_template("project_update", "sample_data")
    print(f"📋 Project update generated: {project_filename}")

    print("\n🎯 PRESENTATION FEATURES DEMONSTRATED:")
    print("- Automated slide creation from data")
    print("- AI-powered content generation")
    print("- Data visualization integration")
    print("- Professional theme application")
    print("- Multiple layout types")
    print("- Chart and graph embedding")

    print("\n🔗 INTEGRATION POSSIBILITIES:")
    print("- Database connectivity for live data")
    print("- API integration for real-time metrics")
    print("- Excel/CSV data import automation")
    print("- Template library management")
    print("- Automated report scheduling")
    print("- Brand guideline enforcement")
    print("- Multi-language support")
    print("- Cloud storage integration")

    print(f"\n💾 Generated files:")
    print(f"- {filename}")
    print(f"- {sales_filename}")
    print(f"- {project_filename}")

if __name__ == "__main__":
    main()